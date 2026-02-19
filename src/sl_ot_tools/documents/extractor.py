"""Per-format text extraction from office documents.

Extracts text content as markdown from .docx, .pptx, .xlsx, and .pdf files
using pure-Python libraries.
"""

from pathlib import Path


def extract_docx(path: Path) -> str:
    """Extract text from a .docx file as markdown.

    Converts heading styles to markdown headings, tables to markdown tables.
    """
    from docx import Document

    doc = Document(str(path))
    lines = []

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            # It's a paragraph — find matching Paragraph object
            for para in doc.paragraphs:
                if para._element is element:
                    text = para.text.strip()
                    if not text:
                        lines.append("")
                        break
                    style = (para.style.name or "").lower()
                    if "heading 1" in style:
                        lines.append(f"# {text}")
                    elif "heading 2" in style:
                        lines.append(f"## {text}")
                    elif "heading 3" in style:
                        lines.append(f"### {text}")
                    elif "heading" in style:
                        lines.append(f"#### {text}")
                    elif "title" in style:
                        lines.append(f"# {text}")
                    else:
                        lines.append(text)
                    break

        elif tag == "tbl":
            # It's a table — find matching Table object
            for table in doc.tables:
                if table._element is element:
                    lines.append(_table_to_markdown(table))
                    break

    return "\n".join(lines)


def _table_to_markdown(table) -> str:
    """Convert a docx table to markdown format."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")

    if len(rows) >= 1:
        # Add header separator after first row
        col_count = len(table.rows[0].cells)
        sep = "| " + " | ".join(["---"] * col_count) + " |"
        rows.insert(1, sep)

    return "\n".join(rows)


def extract_pptx(path: Path) -> str:
    """Extract text from a .pptx file as markdown.

    Outputs slide-by-slide with titles and text frames.
    """
    from pptx import Presentation

    prs = Presentation(str(path))
    lines = []

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")

        # Get title if present
        if slide.shapes.title and slide.shapes.title.text.strip():
            lines.append(f"### {slide.shapes.title.text.strip()}")

        for shape in slide.shapes:
            if shape.has_text_frame:
                # Skip the title shape (already handled)
                if shape == slide.shapes.title:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)

            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    col_count = len(table.columns)
                    sep = "| " + " | ".join(["---"] * col_count) + " |"
                    rows.insert(1, sep)
                    lines.append("\n".join(rows))

        lines.append("")  # blank line between slides

    return "\n".join(lines)


def extract_xlsx(path: Path) -> str:
    """Extract text from an .xlsx file as markdown.

    Outputs sheet-by-sheet with rows as markdown tables.
    """
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    lines = []

    for sheet in wb.sheetnames:
        ws = wb[sheet]
        lines.append(f"## {sheet}")

        rows_data = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell).strip() if cell is not None else "" for cell in row]
            # Skip completely empty rows
            if not any(cells):
                continue
            rows_data.append(cells)

        if not rows_data:
            lines.append("*(empty sheet)*")
            lines.append("")
            continue

        # Normalize column count
        max_cols = max(len(r) for r in rows_data)
        for r in rows_data:
            while len(r) < max_cols:
                r.append("")

        # Build markdown table
        for i, row in enumerate(rows_data):
            escaped = [c.replace("|", "\\|") for c in row]
            lines.append("| " + " | ".join(escaped) + " |")
            if i == 0:
                lines.append("| " + " | ".join(["---"] * max_cols) + " |")

        lines.append("")

    wb.close()
    return "\n".join(lines)


def extract_pdf(path: Path) -> str:
    """Extract text from a .pdf file as markdown.

    Uses pdfplumber for page-by-page text + table extraction.
    """
    import pdfplumber

    lines = []

    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            lines.append(f"## Page {i}")

            # Extract tables first
            tables = page.extract_tables() or []
            if tables:
                for table in tables:
                    if not table:
                        continue
                    rows = []
                    for row in table:
                        cells = [(str(c).strip() if c else "").replace("|", "\\|") for c in row]
                        rows.append("| " + " | ".join(cells) + " |")
                    if rows:
                        col_count = len(table[0]) if table[0] else 1
                        sep = "| " + " | ".join(["---"] * col_count) + " |"
                        rows.insert(1, sep)
                        lines.append("\n".join(rows))
                        lines.append("")

            # Extract remaining text
            text = page.extract_text()
            if text:
                lines.append(text.strip())

            lines.append("")

    return "\n".join(lines)


# Dispatcher mapping file type to extractor function
EXTRACTORS = {
    "docx": extract_docx,
    "pptx": extract_pptx,
    "xlsx": extract_xlsx,
    "pdf": extract_pdf,
}


def extract_text(path: Path) -> str:
    """Extract text from a file based on its extension.

    Args:
        path: Path to the document file.

    Returns:
        Extracted text as markdown string.

    Raises:
        ValueError: If file type is not supported.
    """
    ext = path.suffix.lstrip(".").lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        raise ValueError(f"Unsupported file type: {ext}")
    return extractor(path)
